"""PDF and PPTX export."""

from __future__ import annotations

import os
import shutil
import subprocess
from pathlib import Path


def _find_browser() -> str | None:
    """Find a headless-capable browser on the system."""
    # macOS application paths
    mac_paths = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
    ]

    # Check macOS application paths first
    for path in mac_paths:
        if os.path.isfile(path):
            return path

    # Check PATH for common browser commands
    for cmd in ["google-chrome", "chromium", "chromium-browser", "microsoft-edge", "brave"]:
        found = shutil.which(cmd)
        if found:
            return found

    return None


def _is_chromium_based(browser_path: str) -> bool:
    """Check if the browser is Chromium-based (supports --print-to-pdf)."""
    name = browser_path.lower()
    return any(x in name for x in ["chrome", "chromium", "edge", "brave"])


def export_pdf(input_path: str, output_path: str | None = None) -> str | None:
    """Export a presentation to PDF using a system browser.

    Returns the output path on success, or None if no browser was found.
    """
    from colloquium.build import build_file

    input_path = str(Path(input_path).resolve())

    # Build HTML first
    html_path = build_file(input_path)

    if output_path is None:
        output_path = str(Path(input_path).with_suffix(".pdf"))

    html_url = f"file://{Path(html_path).resolve()}"

    browser = _find_browser()
    if browser is None:
        return None

    if not _is_chromium_based(browser):
        return None

    # Use Chromium's headless print-to-pdf
    # 10in x 5.625in = 16:9 landscape, matching @page size in print CSS
    cmd = [
        browser,
        "--headless",
        "--disable-gpu",
        f"--print-to-pdf={output_path}",
        "--no-pdf-header-footer",
        "--print-to-pdf-no-header",
        "--no-margins",
        f"--paper-width=10",
        f"--paper-height=5.625",
        "--virtual-time-budget=5000",
        html_url,
    ]

    try:
        subprocess.run(cmd, capture_output=True, timeout=30, check=True)
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
        return None

    if Path(output_path).exists():
        # Optional ghostscript compression
        _compress_pdf(output_path)
        return output_path

    return None


def _compress_pdf(pdf_path: str) -> None:
    """Compress PDF with ghostscript if available."""
    gs = shutil.which("gs")
    if not gs:
        return

    compressed = pdf_path + ".compressed"
    cmd = [
        gs,
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        "-dPDFSETTINGS=/prepress",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        f"-sOutputFile={compressed}",
        pdf_path,
    ]

    try:
        subprocess.run(cmd, capture_output=True, timeout=30, check=True)
        # Only replace if compressed is smaller
        if Path(compressed).stat().st_size < Path(pdf_path).stat().st_size:
            os.replace(compressed, pdf_path)
        else:
            os.unlink(compressed)
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError, OSError):
        if os.path.exists(compressed):
            os.unlink(compressed)


# ===== PPTX export =====

import re as _re

_IMAGE_RE = _re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
_CHART_BLOCK_RE = _re.compile(r"```chart\n(.*?)```", _re.DOTALL)
_CODE_BLOCK_RE = _re.compile(r"```\w*\n(.*?)```", _re.DOTALL)
_CONV_BLOCK_RE = _re.compile(r"```conversation\n(.*?)```", _re.DOTALL)
_TABLE_RE = _re.compile(r"((?:\|.+\|\n?)+)", _re.MULTILINE)
_TABLE_SEP_RE = _re.compile(r"^\|[-\s:|]+\|$")


def _strip_markdown(text: str) -> str:
    """Strip markdown formatting to plain text."""
    # Remove images entirely (handled separately)
    text = _re.sub(r"!\[([^\]]*)\]\([^)]*\)", "", text)
    # Remove links, keep text
    text = _re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    # Remove HTML tags
    text = _re.sub(r"<[^>]+>", "", text)
    # Remove bold/italic markers
    text = _re.sub(r"\*\*\*(.+?)\*\*\*", r"\1", text)
    text = _re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = _re.sub(r"\*(.+?)\*", r"\1", text)
    text = _re.sub(r"___(.+?)___", r"\1", text)
    text = _re.sub(r"__(.+?)__", r"\1", text)
    text = _re.sub(r"_(.+?)_", r"\1", text)
    # Remove inline code backticks
    text = _re.sub(r"`(.+?)`", r"\1", text)
    # Remove citation syntax [@key]
    text = _re.sub(r"\[@[^\]]+\]", "", text)
    # Remove HTML comment directives
    text = _re.sub(r"<!--.*?-->", "", text, flags=_re.DOTALL)
    # Remove blockquote markers
    text = _re.sub(r"^>\s?", "", text, flags=_re.MULTILINE)
    # Remove column separators
    text = _re.sub(r"^\|\|\|$", "", text, flags=_re.MULTILINE)
    return text.strip()


def _extract_images(content: str) -> list[tuple[str, str]]:
    """Extract (alt, src) image tuples from markdown content."""
    return _IMAGE_RE.findall(content)


def _extract_chart_specs(content: str) -> list[dict]:
    """Extract chart YAML specs from ```chart blocks."""
    import yaml

    specs = []
    for m in _CHART_BLOCK_RE.finditer(content):
        try:
            spec = yaml.safe_load(m.group(1))
            if isinstance(spec, dict):
                specs.append(spec)
        except Exception:
            pass
    return specs


def _extract_tables(content: str) -> list[list[list[str]]]:
    """Extract markdown tables as list of rows (list of cells)."""
    tables = []
    for m in _TABLE_RE.finditer(content):
        block = m.group(1).strip()
        rows = []
        for line in block.splitlines():
            line = line.strip()
            if not line or _TABLE_SEP_RE.match(line):
                continue
            cells = [c.strip() for c in line.strip("|").split("|")]
            rows.append(cells)
        if len(rows) >= 2:  # at least header + one data row
            tables.append(rows)
    return tables


def _strip_special_blocks(content: str) -> str:
    """Remove images, chart/code/conversation blocks, and tables from content."""
    # Remove fenced blocks (chart, code, conversation)
    text = _re.sub(r"```\w*\n.*?```", "", content, flags=_re.DOTALL)
    # Remove images
    text = _IMAGE_RE.sub("", text)
    # Remove tables
    text = _TABLE_RE.sub("", text)
    return text


def _parse_bullets(content: str) -> list[tuple[int, str]]:
    """Parse cleaned markdown text into (indent_level, text) bullet pairs."""
    items: list[tuple[int, str]] = []
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        m = _re.match(r"^(\s*)([-*+]|\d+\.)\s+(.*)", line)
        if m:
            indent = len(m.group(1))
            level = indent // 2
            items.append((level, _strip_markdown(m.group(3))))
        else:
            plain = _strip_markdown(stripped)
            if plain:
                items.append((0, plain))
    return items


def _resolve_image_path(src: str, base_dir: Path) -> Path | None:
    """Resolve an image src to an absolute path, if it's a local file."""
    if src.startswith(("http://", "https://", "data:")):
        return None  # skip URLs for now
    p = base_dir / src
    if p.exists():
        return p
    return None


def export_pptx(input_path: str, output_path: str | None = None) -> str:
    """Export a presentation to PPTX.

    Requires the `python-pptx` package (install with `pip install colloquium[pptx]`).
    Returns the output path.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.chart.data import CategoryChartData, XyChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError(
            "PPTX export requires python-pptx. "
            "Install it with: pip install colloquium[pptx]"
        )

    from colloquium.parse import parse_file

    input_path = str(Path(input_path).resolve())
    base_dir = Path(input_path).parent
    deck = parse_file(input_path)

    if output_path is None:
        output_path = str(Path(input_path).with_suffix(".pptx"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Content area constants
    C_LEFT = Inches(0.67)
    C_TOP = Inches(1.8)
    C_WIDTH = Inches(12)
    C_HEIGHT = Inches(5.2)

    def _get_placeholder(slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _add_text_bullets(slide, bullets, left, top, width, height):
        """Add a text box with bullet points."""
        if not bullets:
            return
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (level, text) in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = min(level, 4)
            p.font.size = Pt(18)

    def _add_image(slide, img_path, left, top, max_w, max_h):
        """Add an image, scaling to fit within max dimensions."""
        import io
        from PIL import Image as PILImage

        with PILImage.open(img_path) as img:
            img_w, img_h = img.size

            # Convert unsupported formats (WebP, etc.) to PNG in memory
            fmt = img.format or ""
            supported = {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}
            if fmt.upper() not in supported:
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                buf.seek(0)
                img_file = buf
            else:
                img_file = str(img_path)

        aspect = img_w / img_h
        box_aspect = max_w / max_h

        if aspect > box_aspect:
            w = max_w
            h = int(max_w / aspect)
        else:
            h = max_h
            w = int(max_h * aspect)

        x = left + (max_w - w) // 2
        y = top + (max_h - h) // 2

        slide.shapes.add_picture(img_file, x, y, w, h)

    def _add_chart(slide, spec, left, top, width, height):
        """Add a native PPTX chart from a chart YAML spec."""
        chart_type_str = spec.get("type", "bar")
        data = spec.get("data", {})
        title = spec.get("title", "")

        labels = [str(l) for l in data.get("labels", [])]
        datasets = data.get("datasets", [])

        if not datasets:
            return

        type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }
        xl_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        if chart_type_str == "scatter":
            chart_data = XyChartData()
            for ds in datasets:
                series = chart_data.add_series(ds.get("label", "Series"))
                points = ds.get("data", [])
                if points and isinstance(points[0], dict):
                    for point in points:
                        if "x" in point and "y" in point:
                            series.add_data_point(point["x"], point["y"])
                else:
                    for x, y in zip(labels, points):
                        try:
                            series.add_data_point(float(x), float(y))
                        except (TypeError, ValueError):
                            continue
        else:
            if not labels:
                return

            chart_data = CategoryChartData()
            chart_data.categories = labels

            for ds in datasets:
                chart_data.add_series(
                    ds.get("label", "Series"),
                    ds.get("data", []),
                )

        chart_frame = slide.shapes.add_chart(
            xl_type, left, top, width, height, chart_data,
        )
        chart = chart_frame.chart
        chart.has_legend = len(datasets) > 1

        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # Apply colors from spec
        colors = [
            "#0f3460", "#e94560", "#16213e", "#0ea5e9",
            "#10b981", "#f59e0b", "#8b5cf6", "#ec4899",
        ]
        for i, ds in enumerate(datasets):
            color_hex = ds.get("color", colors[i % len(colors)]).lstrip("#")
            try:
                series = chart.series[i]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
                if chart_type_str == "line":
                    series.format.line.color.rgb = RGBColor.from_string(color_hex)
                    series.format.fill.background()
            except Exception:
                pass

    def _add_table(slide, table_data, left, top, width, height):
        """Add a PPTX table from parsed markdown table data."""
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        if rows == 0 or cols == 0:
            return

        row_height = min(height // rows, Inches(0.5))
        tbl_height = row_height * rows

        shape = slide.shapes.add_table(rows, cols, left, top, width, tbl_height)
        table = shape.table

        for r, row_data in enumerate(table_data):
            for c, cell_text in enumerate(row_data):
                if c < cols:
                    cell = table.cell(r, c)
                    cell.text = _strip_markdown(cell_text)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        if r == 0:
                            p.font.bold = True

    def _add_code_block(slide, code, left, top, width, height):
        """Add a code block as a styled text box."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        # Background fill
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

        for i, line in enumerate(code.strip().splitlines()):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.name = "Courier New"

    for slide_data in deck.slides:
        if slide_data.is_title_slide:
            layout = prs.slide_layouts[0]  # Title Slide
            pptx_slide = prs.slides.add_slide(layout)

            if pptx_slide.shapes.title and slide_data.title:
                pptx_slide.shapes.title.text = slide_data.title

            if slide_data.content:
                subtitle_ph = _get_placeholder(pptx_slide, 1)
                if subtitle_ph:
                    subtitle_ph.text = _strip_markdown(slide_data.content)
        else:
            # Use Title Only layout — we'll position content manually
            layout = prs.slide_layouts[5]  # Title Only
            pptx_slide = prs.slides.add_slide(layout)

            if pptx_slide.shapes.title and slide_data.title:
                pptx_slide.shapes.title.text = slide_data.title

            if not slide_data.content:
                if slide_data.speaker_notes:
                    notes_slide = pptx_slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data.speaker_notes
                continue

            content = slide_data.content

            # Extract special content types
            images = _extract_images(content)
            chart_specs = _extract_chart_specs(content)
            tables = _extract_tables(content)
            code_blocks = [m.group(1) for m in _CODE_BLOCK_RE.finditer(content)
                           if not m.group(0).startswith("```chart") and
                           not m.group(0).startswith("```conversation")]

            # Get remaining text after stripping special blocks
            remaining = _strip_special_blocks(content)
            bullets = _parse_bullets(remaining)

            # Layout: divide content area among elements
            elements = []
            if bullets:
                elements.append(("bullets", bullets))
            for img_alt, img_src in images:
                img_path = _resolve_image_path(img_src, base_dir)
                if img_path:
                    elements.append(("image", img_path))
            for spec in chart_specs:
                elements.append(("chart", spec))
            for tbl in tables:
                elements.append(("table", tbl))
            for code in code_blocks:
                elements.append(("code", code))

            if not elements:
                if slide_data.speaker_notes:
                    notes_slide = pptx_slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data.speaker_notes
                continue

            # Single element: give it the full content area
            # For image-only slides, use the full slide for better centering
            if len(elements) == 1:
                kind, data = elements[0]
                if kind == "image" and not slide_data.title:
                    # No title — center image on full slide
                    _add_image(pptx_slide, data, Inches(0.5), Inches(0.5),
                               prs.slide_width - Inches(1), prs.slide_height - Inches(1))
                elif kind == "image":
                    _add_image(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
                elif kind == "bullets":
                    _add_text_bullets(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
                elif kind == "chart":
                    _add_chart(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
                elif kind == "table":
                    _add_table(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
                elif kind == "code":
                    _add_code_block(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
            else:
                # Multiple elements: split horizontally (text left, visual right)
                # or stack vertically
                has_visual = any(k in ("image", "chart") for k, _ in elements)
                text_els = [(k, d) for k, d in elements if k in ("bullets", "code")]
                visual_els = [(k, d) for k, d in elements if k in ("image", "chart")]
                table_els = [(k, d) for k, d in elements if k == "table"]

                if has_visual and text_els:
                    # Side by side: text left, visual right
                    half_w = C_WIDTH // 2 - Inches(0.2)
                    # Text on left
                    y_offset = C_TOP
                    for kind, data in text_els:
                        h = C_HEIGHT // len(text_els) if len(text_els) > 1 else C_HEIGHT
                        if kind == "bullets":
                            _add_text_bullets(pptx_slide, data, C_LEFT, y_offset, half_w, h)
                        elif kind == "code":
                            _add_code_block(pptx_slide, data, C_LEFT, y_offset, half_w, h)
                        y_offset += h

                    # Visuals on right
                    right_left = C_LEFT + half_w + Inches(0.4)
                    y_offset = C_TOP
                    vis_h = C_HEIGHT // len(visual_els) if visual_els else C_HEIGHT
                    for kind, data in visual_els:
                        if kind == "image":
                            _add_image(pptx_slide, data, right_left, y_offset, half_w, vis_h)
                        elif kind == "chart":
                            _add_chart(pptx_slide, data, right_left, y_offset, half_w, vis_h)
                        y_offset += vis_h

                    # Tables below
                    for kind, data in table_els:
                        _add_table(pptx_slide, data, C_LEFT, C_TOP, C_WIDTH, C_HEIGHT)
                else:
                    # Stack vertically
                    n = len(elements)
                    slot_h = C_HEIGHT // n
                    y_offset = C_TOP
                    for kind, data in elements:
                        if kind == "bullets":
                            _add_text_bullets(pptx_slide, data, C_LEFT, y_offset, C_WIDTH, slot_h)
                        elif kind == "image":
                            _add_image(pptx_slide, data, C_LEFT, y_offset, C_WIDTH, slot_h)
                        elif kind == "chart":
                            _add_chart(pptx_slide, data, C_LEFT, y_offset, C_WIDTH, slot_h)
                        elif kind == "table":
                            _add_table(pptx_slide, data, C_LEFT, y_offset, C_WIDTH, slot_h)
                        elif kind == "code":
                            _add_code_block(pptx_slide, data, C_LEFT, y_offset, C_WIDTH, slot_h)
                        y_offset += slot_h

        # Speaker notes
        if slide_data.speaker_notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

    prs.save(output_path)
    return output_path
